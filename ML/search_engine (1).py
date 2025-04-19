#!/usr/bin/env python
# coding: utf-8

# In[87]:


import os
import papermill as pm
import faiss
import pickle
import numpy as np
from sentence_transformers import SentenceTransformer
from pptx import Presentation


# Функция для проверки обновлений данных и перестройки индекса
def check_updates():
    # Путь к новому JSON с презентациями
    data_time = os.path.getmtime('data/notes.json')  # Обновленный файл с данными
    index_time = os.path.getmtime('index/faiss.index') if os.path.exists('index/faiss.index') else 0
    
    if data_time > index_time:
     #  print("Обнаружены новые данные, перестраиваю индекс...")
        # Запускаем build_index.ipynb, чтобы создать новый индекс
        pm.execute_notebook('bulild_index.ipynb', 'bulild_index_output.ipynb')
   # else:
    #    print('Обновлений нет')


# Проверяем наличие обновлений при каждом запуске
check_updates()


# Загружаем модель для преобразования текста в эмбеддинги
model = SentenceTransformer('all-MiniLM-L6-v2')

# Загружаем индекс FAISS и метаданные (JSON о презентациях)
index = faiss.read_index("index/faiss.index")
with open("index/meta.pkl", "rb") as f:
    documents = pickle.load(f)


# Функция для извлечения текста из презентаций (файлов .pptx)
def extract_text_from_pptx(file_name):
    presentation = Presentation(f"data/files/{file_name}")
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


# In[88]:


# Функция для поиска по индексу
def search(query, top_k=3):
    query_embedding = model.encode([query])  # Преобразуем запрос в эмбеддинг
    
    # Выполняем поиск по индексу
    D, I = index.search(np.array(query_embedding), top_k)
    
    results = []
    for idx, distance in zip(I[0], D[0]):  # Итерируем через индексы и соответствующие расстояния
        similarity = (1 - distance) * 100  # Вычисляем схожесть как процент
        results.append({
            "index": idx,
            "distance": distance,
            "similarity": similarity,  # Добавляем процент схожести
            "document": documents[idx]
        })
    
    # Сортируем результаты по расстоянию (от меньшего к большему), то есть по схожести
    results.sort(key=lambda x: x["distance"])  # Чем меньше distance, тем выше схожесть
    
    # Возвращаем отсортированные документы с полем 'similarity'
    return results  # Возвращаем весь список, а не только документы

# Функция для форматирования ответа (с использованием данных о преподавателе, предмете и презентациях)
def format_answer(results):
    formatted_results = []
    for result in results:
        # Получаем данные из документа
        subject = result["document"]["subject"]
        title = result["document"]["title"]
        teacher_name = result["document"]["teacher"]
        file_name = result["document"]["file_name"]
        similarity = result["similarity"]  # Процент схожести
        
        # Получаем текст презентации (если нужно)
        pptx_text = extract_text_from_pptx(f"{file_name}")
        
        formatted_results.append({
            "subject": subject,
            "title": title,
            "teacher": teacher_name,
            "file_name": file_name,
            "pptx_text": pptx_text,  # Можно добавить текст конспекта здесь
            "similarity": similarity  # Добавляем процент схожести
        })
    return formatted_results

# Пример запроса (вопрос по теории)
def prompt(query):
    # Ищем по запросу
    results = search(query)
    
    # Форматируем ответ
    formatted_results = format_answer(results)
    
    for result in formatted_results:
        print(f"Предмет: {result['subject']}, Тема: {result['title']}")
        print(f"Преподаватель: {result['teacher']}")
        #print(f"Процент схожести: {result['similarity']:.2f}%")  # Выводим процент схожести
        print("-" * 80)

