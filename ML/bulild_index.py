#!/usr/bin/env python
# coding: utf-8

# In[48]:


import json
from pptx import Presentation

# Загружаем данные
with open('data/notes.json') as file:
    data = json.load(file)

# Создаем словари для пользователей и предметов
users = {u['id']: u for u in data['users']}
subjects = {s['id']: s for s in data['subjects']}

# Функция для извлечения текста из презентации
def extract_text_from_pptx(filename):
    prs = Presentation(f"data/files/{filename}")
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Список документов
documents = []

# Процессируем презентации
for pres in data["presentations"]:
    subject = subjects.get(pres["subjectId"], {})
    teacher = users.get(pres["authorId"], {})
    
    # Извлекаем текст файла
    file_text = extract_text_from_pptx(pres["file"]["name"])
    
    # Объединяем все важные поля
    full_text = f"{pres['title']} {pres['description']} {subject.get('name', '')} {teacher.get('name', '')} {file_text}"
    
    # Создаем запись для документа с нужными данными
    documents.append({
        "id": pres["id"],
        "title": pres["title"],  # Название презентации
        "subject": subject.get("name", ""),  # Название предмета
        "teacher": teacher.get("name", ""),  # Имя преподавателя
        "file_name": pres["file"]["name"],  # Название файла
        "text": full_text  })


# In[49]:


from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import pickle

model = SentenceTransformer('all-MiniLM-L6-v2')

texts = [doc["text"] for doc in documents]
embeddings = model.encode(texts, show_progress_bar=True)

# FAISS индекс
index = faiss.IndexFlatL2(embeddings[0].shape[0])
index.add(np.array(embeddings))

# Сохраняем
faiss.write_index(index, "index/faiss.index")
with open("index/meta.pkl", "wb") as f:
    pickle.dump(documents, f)

